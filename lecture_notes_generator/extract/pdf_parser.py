import pdfplumber
from pptx import Presentation
import os

def parse_pdf_slides(file_path: str) -> str:
    \"\"\"
    Extracts text from a PDF slide presentation.
    Preserves slide boundaries.
    \"\"\"
    text_content = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    text_content.append(f\"--- Slide {i + 1} ---\\n{text.strip()}\\n\")
    except Exception as e:
        print(f\"Error parsing PDF slides {file_path}: {e}\")
    
    return \"\\n\".join(text_content)

def parse_pptx_slides(file_path: str) -> str:
    \"\"\"
    Extracts text from a PPTX presentation.
    Preserves slide boundaries.
    \"\"\"
    text_content = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, \"text\"):
                    slide_text.append(shape.text)
            
            if slide_text:
                full_text = \"\\n\".join(slide_text).strip()
                text_content.append(f\"--- Slide {i + 1} ---\\n{full_text}\\n\")
    except Exception as e:
        print(f\"Error parsing PPTX slides {file_path}: {e}\")
        
    return \"\\n\".join(text_content)

def extract_slides(file_path: str) -> str:
    \"\"\"
    Wrapper function to extract text from slides based on file extension.
    \"\"\"
    ext = os.path.splitext(file_path)[1].lower()
    if ext == \".pdf\":
        return parse_pdf_slides(file_path)
    elif ext == \".pptx\":
        return parse_pptx_slides(file_path)
    else:
        raise ValueError(f\"Unsupported slide format: {ext}\")

def extract_textbook(file_path: str) -> str:
    \"\"\"
    Extracts text from a textbook (PDF).
    Does not preserve slide-like boundaries, just returns raw text or pages.
    \"\"\"
    text_content = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_content.append(text.strip())
    except Exception as e:
        print(f\"Error parsing PDF textbook {file_path}: {e}\")
    
    return \"\\n\\n\".join(text_content)
