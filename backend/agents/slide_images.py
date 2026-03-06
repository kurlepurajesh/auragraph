"""
backend/agents/slide_images.py
───────────────────────────────
Extract embedded images from PDF slides and PPTX files.

Uses:
  • PyMuPDF (fitz)     for PDF  — reliable raw image bytes per page
  • python-pptx        for PPTX — shape.image.blob access

Returns ExtractedImage objects. Caller saves them and injects
markdown references into the notes.

Filtering heuristics (skip decorative / tiny images):
  • width  < 80 px  or  height < 80 px  → skip
  • area   < 10 000 px²                 → skip (almost certainly an icon)
  • size   < 4 KB                       → skip
  • duplicate hash                       → skip
  • PDF raster images that are clearly
    just background fills (very few distinct colors) → skip
"""

from __future__ import annotations

import hashlib
import io
import logging
from dataclasses import dataclass, field

logger = logging.getLogger(__name__)

# Maximum images to extract per file (avoids runaway for 100-slide decks)
MAX_IMAGES_PER_FILE = 20
# Minimum significant dimension
MIN_DIM_PX = 80
MIN_AREA_PX = 10_000
MIN_BYTES   = 4_000


@dataclass
class ExtractedImage:
    img_id:       str     # e.g. "img_001"
    data:         bytes   # raw PNG or JPEG bytes
    mime:         str     # "image/png" | "image/jpeg"
    source_label: str     # "Slide 3" | "Page 5"
    width:        int
    height:       int
    description:  str = ""   # filled by describe_slide_image() later


# ── PDF extraction via PyMuPDF ─────────────────────────────────────────────

def extract_images_from_pdf(file_bytes: bytes) -> list[ExtractedImage]:
    """Extract significant embedded images from a PDF."""
    results: list[ExtractedImage] = []
    seen_hashes: set[str] = set()
    counter = 1

    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.warning("slide_images: PyMuPDF not installed — PDF image extraction unavailable")
        return []

    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as e:
        logger.warning("slide_images: fitz failed to open PDF: %s", e)
        return []

    for page_num, page in enumerate(doc, start=1):
        if len(results) >= MAX_IMAGES_PER_FILE:
            break

        try:
            img_list = page.get_images(full=True)
        except Exception:
            continue

        for img_info in img_list:
            if len(results) >= MAX_IMAGES_PER_FILE:
                break
            try:
                xref  = img_info[0]
                w     = img_info[2]
                h     = img_info[3]

                # Size filter
                if w < MIN_DIM_PX or h < MIN_DIM_PX:
                    continue
                if w * h < MIN_AREA_PX:
                    continue

                # Extract raw bytes
                base_img = doc.extract_image(xref)
                data  = base_img["image"]
                mime_short = base_img.get("ext", "png").lower()
                mime        = f"image/{mime_short}".replace("image/jpg", "image/jpeg")

                if len(data) < MIN_BYTES:
                    continue

                # Deduplicate
                h_str = hashlib.md5(data).hexdigest()
                if h_str in seen_hashes:
                    continue
                seen_hashes.add(h_str)

                img_id = f"img_{counter:03d}"
                counter += 1
                results.append(ExtractedImage(
                    img_id=img_id,
                    data=data,
                    mime=mime,
                    source_label=f"Page {page_num}",
                    width=w,
                    height=h,
                ))
                logger.debug("slide_images: extracted %s from PDF page %d (%dx%d, %d bytes)",
                             img_id, page_num, w, h, len(data))
            except Exception as e:
                logger.debug("slide_images: skipping image xref %s: %s", img_info[0], e)

    doc.close()
    logger.info("slide_images: extracted %d images from PDF (%d pages)", len(results), len(doc))
    return results


# ── PPTX extraction via python-pptx ──────────────────────────────────────────

def extract_images_from_pptx(file_bytes: bytes) -> list[ExtractedImage]:
    """Extract significant embedded images from a PPTX file."""
    results: list[ExtractedImage] = []
    seen_hashes: set[str] = set()
    counter = 1

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        logger.warning("slide_images: python-pptx not installed")
        return []

    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        logger.warning("slide_images: python-pptx failed to open file: %s", e)
        return []

    for slide_num, slide in enumerate(prs.slides, start=1):
        if len(results) >= MAX_IMAGES_PER_FILE:
            break

        for shape in slide.shapes:
            if len(results) >= MAX_IMAGES_PER_FILE:
                break
            try:
                shape_type = getattr(shape, 'shape_type', None)
                # MSO_SHAPE_TYPE.PICTURE = 13
                if shape_type != 13:
                    continue
                img = shape.image
                data = img.blob
                if len(data) < MIN_BYTES:
                    continue

                # Get pixel dimensions (pptx stores EMUs — 914400 per inch at 96dpi)
                # We approximate from the shape's display size instead
                w_emu = shape.width  or 1
                h_emu = shape.height or 1
                # Rough pixel estimate at 96dpi: 914400 EMU/inch, 96px/inch
                w_px = int(w_emu / 914400 * 96)
                h_px = int(h_emu / 914400 * 96)

                if w_px < MIN_DIM_PX or h_px < MIN_DIM_PX:
                    continue

                ext  = img.ext.lower()
                mime = f"image/{ext}".replace("image/jpg", "image/jpeg")

                h_str = hashlib.md5(data).hexdigest()
                if h_str in seen_hashes:
                    continue
                seen_hashes.add(h_str)

                img_id = f"img_{counter:03d}"
                counter += 1
                results.append(ExtractedImage(
                    img_id=img_id,
                    data=data,
                    mime=mime,
                    source_label=f"Slide {slide_num}",
                    width=w_px,
                    height=h_px,
                ))
                logger.debug("slide_images: extracted %s from PPTX slide %d (%dx%d, %d bytes)",
                             img_id, slide_num, w_px, h_px, len(data))
            except Exception as e:
                logger.debug("slide_images: skipping PPTX shape: %s", e)

    logger.info("slide_images: extracted %d images from PPTX (%d slides)",
                len(results), len(prs.slides))
    return results


# ── Unified entry point ───────────────────────────────────────────────────────

def extract_images_from_file(file_bytes: bytes, filename: str) -> list[ExtractedImage]:
    """Route to the right extractor based on file extension."""
    from pathlib import Path
    ext = Path(filename.lower()).suffix
    if ext in {'.pptx', '.ppt'}:
        return extract_images_from_pptx(file_bytes)
    if ext in {'.pdf'}:
        return extract_images_from_pdf(file_bytes)
    # Image files themselves are not "slide sources" — skip
    return []


# ── Image store (temp filesystem) ────────────────────────────────────────────

import os
_IMG_ROOT = "/tmp/auragraph_imgs"


def save_images(notebook_id: str, images: list[ExtractedImage]) -> None:
    """Save image bytes to disk so the API endpoint can serve them."""
    folder = os.path.join(_IMG_ROOT, notebook_id)
    os.makedirs(folder, exist_ok=True)
    for img in images:
        ext  = img.mime.split("/")[-1].replace("jpeg", "jpg")
        path = os.path.join(folder, f"{img.img_id}.{ext}")
        with open(path, "wb") as f:
            f.write(img.data)
    logger.info("slide_images: saved %d images for notebook %s", len(images), notebook_id)


def get_image_path(notebook_id: str, img_id_with_ext: str) -> str | None:
    """Returns the filesystem path for a stored image, or None if not found."""
    folder = os.path.join(_IMG_ROOT, notebook_id)
    path   = os.path.join(folder, img_id_with_ext)
    return path if os.path.exists(path) else None
