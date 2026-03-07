"""
PDF / PPTX Extraction Utility  (RAGFlow-inspired)
--------------------------------------------------
Key improvements over the previous PyPDF2-only version:
  • pdfplumber for PDF — preserves reading order and avoids the
    column-merging artefacts that PyPDF2 produces on two-column slides.
  • python-pptx for native PPTX — extracts slide text in reading order
    (top-to-bottom, left-to-right), preserving bullet hierarchy.
  • Slide-boundary-aware chunking — each slide / page becomes its own
    unit so the summariser can see topic boundaries instead of a wall of text.
  • Proportional budget allocation in summarise_chunks() ensures every
    topic is represented even when the combined text exceeds the LLM window.
"""

from __future__ import annotations

import io
import logging
import re

logger = logging.getLogger(__name__)

# ──────────────────────────────────────────────────────────────────────────────
# PDF extraction artifact scrubber
# ──────────────────────────────────────────────────────────────────────────────

_CID_RE          = re.compile(r'\(cid:\d+\)')
_PAGE_MARKER_RE  = re.compile(r'---\s*[Pp]age\s+\d+\s*---')
_DASH_ONLY_RE    = re.compile(r'^[\s\u2212\u2013\u2014\-]+$')


def _scrub_pdf_artifacts(text: str) -> str:
    """
    Remove common PDF extraction artifacts:
      1. (cid:N)  — unresolved glyph references from CID-keyed fonts
      2. Single- or two-character lines — glyph-per-line dumps from Type3 fonts
      3. --- Page N --- markers (both compact and spaced variants)
      4. Pure-dash/minus separator lines
      5. Mirror/echo lines — drop the spaced-glyph version when the very next
         non-empty line is the concatenated version of the same characters
    """
    # Pass 1: strip (cid:N) inline
    text = _CID_RE.sub('', text)

    # Pass 2: remove embedded page markers
    text = _PAGE_MARKER_RE.sub('', text)

    # Pass 3: per-line cleanup
    lines = text.split('\n')
    cleaned: list[str] = []
    i = 0
    while i < len(lines):
        raw    = lines[i]
        s      = raw.strip()

        # Pure dash/separator lines
        if s and _DASH_ONLY_RE.match(s):
            i += 1
            continue

        # Single or two character lines (glyph dump fragments).
        # Allow truly empty lines (paragraph breaks).
        if 1 <= len(s) <= 2:
            i += 1
            continue

        # Mirror-text detection: if the next non-empty line is the same text
        # with all spaces removed, skip THIS line (the spaced glyph version)
        # and keep the next (the clean ToUnicode reconstruction).
        if s:
            compact_s = s.replace(' ', '')
            # Look for the next non-empty line
            j = i + 1
            while j < len(lines) and not lines[j].strip():
                j += 1
            if j < len(lines):
                compact_next = lines[j].strip().replace(' ', '')
                if compact_next and compact_next == compact_s and len(compact_s) > 3:
                    # Current line is just the next line with spaces inserted — skip it
                    i += 1
                    continue

        cleaned.append(raw)
        i += 1

    # Pass 4: collapse runs of 3+ blank lines → 2 blank lines
    result = re.sub(r'\n{3,}', '\n\n', '\n'.join(cleaned))
    return result.strip()


# ──────────────────────────────────────────────────────────────────────────────
# PDF extraction (pdfplumber)
# ──────────────────────────────────────────────────────────────────────────────

def _is_cover_page(text: str, page_num: int) -> bool:
    """
    Heuristic: detect cover / title pages that should be skipped.
    Cover pages typically have very few lines of content (< 5 meaningful lines)
    and contain keywords like "university", "department", "course", "semester",
    "instructor", "presented by", "submitted to", or are just the book/course title.
    Only applies to the first 2 pages.
    """
    if page_num > 2:
        return False
    lines = [l.strip() for l in text.strip().split('\n') if l.strip()]
    if len(lines) < 3:
        return True  # near-empty page — definitely skip
    if len(lines) > 10:
        return False  # too much content to be a cover
    text_lower = text.lower()
    cover_signals = [
        'university', 'department of', 'course:', 'course no', 'semester',
        'instructor:', 'professor:', 'presented by', 'submitted to',
        'lecture notes', 'study material', 'prepared by', 'module',
        'unit -', 'unit–', 'subject code', 'class notes'
    ]
    signal_count = sum(1 for s in cover_signals if s in text_lower)
    # If 2+ cover signals and very few content lines — it's a cover/title page
    return signal_count >= 2 and len(lines) <= 8



    """
    Extract full text from a PDF file.

    Strategy (RAGFlow-inspired):
      1. Try pdfplumber first — it uses pdfminer under the hood and correctly
         handles multi-column layouts, preserving reading order.
      2. Fall back to PyPDF2 if pdfplumber fails (e.g. heavily encrypted PDFs).
    Returns one string with pages separated by '\\n\\n--- Page N ---\\n\\n'.
    """
    try:
        import pdfplumber
        pages_text: list[str] = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = page.extract_text(x_tolerance=3, y_tolerance=3)
                if text and text.strip():
                    page_text = text.strip()
                    if _is_cover_page(page_text, i):
                        logger.info("Skipping cover/title page %d", i)
                        continue
                    pages_text.append(f"--- Page {i} ---\n{page_text}")
        if pages_text:
            return _scrub_pdf_artifacts("\n\n".join(pages_text))
    except Exception as e:
        logger.warning("pdfplumber failed (%s), falling back to PyPDF2", e)

    # Fallback: PyPDF2
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        pages_text = []
        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text()
            if text and text.strip():
                pages_text.append(f"--- Page {i} ---\n{text.strip()}")
        if pages_text:
            return _scrub_pdf_artifacts("\n\n".join(pages_text))
    except Exception as e:
        raise ValueError(f"Failed to parse PDF: {e}") from e

    # OCR fallback: PDF is image-based (scanned / handwritten) — render pages and OCR
    logger.info("PDF yielded no selectable text — attempting OCR via image render")
    try:
        ocr_text = _extract_pdf_with_ocr(file_bytes)
        if ocr_text.strip():
            return ocr_text
    except Exception as e:
        logger.warning("PDF OCR fallback failed: %s", e)

    raise ValueError("No text could be extracted from the PDF.")


def _extract_pdf_with_ocr(file_bytes: bytes) -> str:
    """
    Render each page of a scanned / handwritten PDF as an image via PyMuPDF
    and OCR it using the existing image_ocr pipeline (Groq vision → Tesseract).
    Returns page-boundary-marked text identical to the text extraction path.
    """
    import fitz  # PyMuPDF — already in requirements.txt
    from agents.image_ocr import extract_text_from_image

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages_text: list[str] = []
    try:
        for i, page in enumerate(doc, start=1):
            # Render at 2× scale (~144 dpi) for good OCR quality with manageable size
            mat  = fitz.Matrix(2.0, 2.0)
            pix  = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes("png")
            try:
                text = extract_text_from_image(img_bytes, f"page_{i}.png")
            except Exception as e:
                logger.warning("OCR failed on page %d: %s", i, e)
                text = ""
            if text.strip():
                pages_text.append(f"--- Page {i} ---\n{text.strip()}")
    finally:
        doc.close()
    return "\n\n".join(pages_text) if pages_text else ""


# ──────────────────────────────────────────────────────────────────────────────
# PPTX extraction  (RAGFlow ppt_parser.py approach)
# ──────────────────────────────────────────────────────────────────────────────

def _sort_shapes(shapes):
    """
    Sort shapes top-to-bottom, left-to-right (RAGFlow's shape-order heuristic).
    Shapes without position attributes are placed last.
    """
    def _key(s):
        top  = s.top  if s.top  is not None else 9_999_999
        left = s.left if s.left is not None else 9_999_999
        return (top, left)
    return sorted(shapes, key=_key)


def _extract_shape_text(shape) -> str:
    """Recursively extract text from a PPTX shape, preserving bullet hierarchy."""
    try:
        from pptx.util import Pt
        # Text frames (most common)
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            lines: list[str] = []
            for para in shape.text_frame.paragraphs:
                raw = para.text.strip()
                if not raw:
                    continue
                # Detect bullet / numbered list
                is_bullet = (
                    bool(para._p.xpath("./a:pPr/a:buChar"))
                    or bool(para._p.xpath("./a:pPr/a:buAutoNum"))
                    or bool(para._p.xpath("./a:pPr/a:buBlip"))
                )
                indent = "  " * max(0, para.level)
                prefix = "• " if is_bullet else ""
                lines.append(f"{indent}{prefix}{raw}")
            return "\n".join(lines)

        # Tables
        shape_type = None
        try:
            shape_type = shape.shape_type
        except Exception:
            pass

        if shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
            rows = []
            for row in shape.table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    rows.append(" | ".join(cells))
            return "\n".join(rows)

        # Group shapes — recurse
        if shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            texts = [_extract_shape_text(s) for s in _sort_shapes(shape.shapes)]
            return "\n".join(t for t in texts if t)

    except Exception as e:
        logger.debug("Shape extraction error: %s", e)
    return ""


# Regex to detect lines that are metadata, not teaching content.
# Rules:
#   - email addresses
#   - lines starting with an honorific (Dr/Prof/Mr/Mrs/Ms) + name(s)
#   - numeric dates: dd/mm/yyyy, dd-mm-yyyy
#   - month + year: May 2024
# Deliberately does NOT match bare "First Last" to avoid stripping topic names
# like "Fourier Transform", "Convolution Theorem", etc.
_METADATA_LINE_RE = re.compile(
    r'^\s*('
    r'[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}'           # email
    r'|(?:Dr|Prof|Mr|Mrs|Ms|Er)\.?\s+[A-Z][a-z]+(\s[A-Z][a-z]+)*'  # Dr./Prof. A B C
    r'|\d{1,4}[-\/]\d{1,2}[-\/]\d{2,4}'                            # 01/03/2026
    r'|(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\.?\s+\d{4}'  # May 2024
    r')\s*$',
    re.IGNORECASE,
)

# Lines that contain an institution keyword are almost always metadata
_INSTITUTION_RE = re.compile(
    r'\b(?:Department|School|College|Institute|University|Faculty|IIT|NIT|BITS)\b',
    re.IGNORECASE,
)


def _strip_metadata_lines(text: str) -> str:
    """
    Remove lines that are almost certainly author/institution metadata
    (emails, honorifics, plain person names, dates).
    Only strips lines that match metadata patterns exactly — mathematical
    or sentence-like lines are always kept.
    """
    kept: list[str] = []
    for line in text.splitlines():
        stripped = line.strip()
        # Keep empty lines and any line with teaching content indicators
        if not stripped:
            kept.append(line)
            continue
        # Keep lines that have likely teaching content.
        # NOTE: '/' intentionally excluded from has_math — dates also contain '/'
        has_colon   = ':' in stripped
        has_math    = any(c in stripped for c in ['=', '+', '$', '\\', '\u222b', '\u03a3'])
        has_bullet  = stripped.startswith(('\u2022', '-', '*', '\u2013', '1.', '2.', '3.'))
        is_sentence = len(stripped.split()) > 6
        is_institution = bool(_INSTITUTION_RE.search(stripped))
        if has_colon or has_math or has_bullet or is_sentence:
            kept.append(line)
        elif _METADATA_LINE_RE.match(stripped) or is_institution:
            continue  # drop metadata
        else:
            kept.append(line)
    return '\n'.join(kept)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """
    Extract text from a PPTX file with slide-boundary markers.

    Each slide becomes:
        --- Slide N: <Title> ---
        <body text in reading order>

    This mirrors RAGFlow's RAGFlowPptParser.__call__ but adds title detection.
    """
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        slides_text: list[str] = []

        for i, slide in enumerate(prs.slides, start=1):
            # Try to get slide title
            title_text = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text.strip()

            body_parts: list[str] = []
            for shape in _sort_shapes(slide.shapes):
                # Skip the title shape — it's already captured
                if slide.shapes.title and shape == slide.shapes.title:
                    continue
                txt = _extract_shape_text(shape)
                if txt:
                    body_parts.append(txt)

            body = _strip_metadata_lines("\n".join(body_parts)).strip()

            # Build slide block
            if title_text:
                header = f"--- Slide {i}: {title_text} ---"
            else:
                header = f"--- Slide {i} ---"

            if body or title_text:
                content = f"{header}\n{body}" if body else header
                slides_text.append(content)

        return "\n\n".join(slides_text) if slides_text else ""
    except Exception as e:
        raise ValueError(f"Failed to parse PPTX: {e}") from e


def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    """
    Unified entry point: routes to the right extractor based on file extension.
    Supports: .pdf, .pptx/.ppt, .jpg/.jpeg/.png/.webp/.bmp/.tiff/.heic/.heif
    """
    from pathlib import Path
    fname = filename.lower()
    if fname.endswith(".pptx") or fname.endswith(".ppt"):
        return extract_text_from_pptx(file_bytes)
    # Image formats → OCR (Groq vision → tesseract → placeholder)
    if Path(fname).suffix in {
        '.jpg', '.jpeg', '.png', '.webp',
        '.bmp', '.tiff', '.tif', '.heic', '.heif',
    }:
        from agents.image_ocr import extract_text_from_image
        return extract_text_from_image(file_bytes, filename)
    return extract_text_from_pdf(file_bytes)


# ──────────────────────────────────────────────────────────────────────────────
# Slide-boundary-aware chunking
# ──────────────────────────────────────────────────────────────────────────────

_SLIDE_BOUNDARY = re.compile(r"^--- Slide \d+", re.MULTILINE)
_PAGE_BOUNDARY  = re.compile(r"^--- Page \d+",  re.MULTILINE)


def chunk_text(text: str, max_chars: int = 8000) -> list[str]:
    """
    Split text into chunks, respecting slide/page boundaries first.

    RAGFlow insight: each slide is a self-contained teaching unit — don't
    split mid-slide.  We group whole slides into chunks up to max_chars,
    then fall back to paragraph splitting for oversized individual slides.
    """
    # Detect whether the text has slide or page markers
    has_slides = bool(_SLIDE_BOUNDARY.search(text))
    has_pages  = bool(_PAGE_BOUNDARY.search(text))
    boundary   = _SLIDE_BOUNDARY if has_slides else _PAGE_BOUNDARY

    if has_slides or has_pages:
        # Split on the boundary markers, keep the marker with its content
        parts = re.split(r"(?=^--- (?:Slide|Page) )", text, flags=re.MULTILINE)
        parts = [p.strip() for p in parts if p.strip()]
    else:
        # Plain paragraph splitting
        parts = [p.strip() for p in text.split("\n\n") if p.strip()]

    chunks: list[str] = []
    current = ""

    for part in parts:
        if len(current) + len(part) + 2 <= max_chars:
            current += ("\n\n" if current else "") + part
        else:
            if current:
                chunks.append(current)
            # Part itself is larger than max_chars — split at paragraph boundary
            if len(part) > max_chars:
                sub_paras = part.split("\n\n")
                sub_buf = ""
                for sp in sub_paras:
                    if len(sub_buf) + len(sp) + 2 <= max_chars:
                        sub_buf += ("\n\n" if sub_buf else "") + sp
                    else:
                        if sub_buf:
                            chunks.append(sub_buf)
                        # FIX (round 4): a single paragraph larger than max_chars
                        # must itself be hard-split rather than stored as-is,
                        # otherwise oversized chunks can exceed LLM context limits.
                        if len(sp) > max_chars:
                            for i in range(0, len(sp), max_chars):
                                chunks.append(sp[i:i + max_chars])
                            sub_buf = ""
                        else:
                            sub_buf = sp
                if sub_buf:
                    current = sub_buf
                else:
                    current = ""
            else:
                current = part

    if current:
        chunks.append(current)

    return chunks or [text]


def summarise_chunks(chunks: list[str], max_summary_chars: int = 32000) -> str:
    """
    Collapse chunks into one string up to max_summary_chars.

    Design principles:
    1. GPT-4o has a 128k token context. 32k chars ≈ 8k tokens — well within
       budget for typical lecture material. The old default of 6k-10k chars
       was far too aggressive and starved the model of content.
    2. Proportional sampling: every chunk (slide/page) gets a fair share of
       the budget so no topic is completely dropped.
    3. Sentence-boundary trimming: never cut mid-sentence; find the last
       sentence-ending punctuation within the budget window.
    4. Slide markers are always preserved (never trimmed away) so GPT-4o
       can see the section structure even in a trimmed chunk.
    """
    combined = "\n\n---\n\n".join(chunks)
    if len(combined) <= max_summary_chars:
        return combined

    # Each chunk gets a proportional share of the budget.
    # Floor at 500 chars so every slide gets at least its title + first sentence.
    per_chunk = max(500, max_summary_chars // max(len(chunks), 1))
    parts: list[str] = []

    for chunk in chunks:
        if len(chunk) <= per_chunk:
            parts.append(chunk)
            continue

        # Always preserve the slide/page header line (first line of the chunk)
        lines = chunk.split("\n")
        header = lines[0] if lines and lines[0].startswith("---") else ""
        body   = chunk[len(header):].strip() if header else chunk

        budget = per_chunk - len(header) - 2  # 2 for "\n\n"
        if budget <= 0:
            parts.append(header)
            continue

        trimmed = body[:budget]

        # Find the last clean sentence boundary: ". ", ".\n", or "\n\n"
        last_stop = max(
            trimmed.rfind(". "),
            trimmed.rfind(".\n"),
            trimmed.rfind("\n\n"),
            trimmed.rfind("! "),
            trimmed.rfind("? "),
        )
        # Only use the boundary if it's in the second half of the window
        # (avoids cutting to just the first sentence of a long section)
        if last_stop > budget // 3:
            trimmed = trimmed[: last_stop + 1].rstrip()

        if header:
            parts.append(f"{header}\n\n{trimmed}")
        else:
            parts.append(trimmed)

    return "\n\n---\n\n".join(parts)
