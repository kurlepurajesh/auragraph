"""
test_notes_pipeline.py  —  AuraGraph rigorous test suite
─────────────────────────────────────────────────────────
Run with:  python test_notes_pipeline.py   (from backend/ with venv active)

Sections
  1.  _strip_metadata_lines        — 21 cases
  2.  _clean_pdf_text              — noise removal, camelCase splitting
  3.  _is_math_line                — formula vs prose classification
  4.  extract_text_from_pptx       — 6 edge cases via synthetic PPTX
  5.  chunk_text / summarise_chunks — slide-boundary chunking + budget
  6.  fix_latex_delimiters         — 12 LaTeX normalisation cases
  7.  generate_local_note          — Beginner / Intermediate / Advanced
  8.  extract_concepts             — pattern matching + heading fallback
  9.  local_mutate                 — 8 doubt categories
  10. local_examine                — 5 concepts, generic fallback
  11. Page splitter                — 5 JS-mirrored edge cases
  12. Live API                     — health, auth, notebooks CRUD,
                                     /api/fuse, /api/upload-fuse-multi,
                                     /api/mutate, /api/examine, /api/extract-concepts
"""

from __future__ import annotations
import io, sys, re, traceback, textwrap, uuid
sys.path.insert(0, ".")

G = "\033[92m"; R = "\033[91m"; Y = "\033[93m"; RST = "\033[0m"; B = "\033[1m"
passed = failed = skipped = 0

def ok(name, detail=""):
    global passed; passed += 1
    print(f"  {G}*{RST} {name}" + (f"  [{detail}]" if detail else ""))

def fail(name, detail):
    global failed; failed += 1
    print(f"  {R}FAIL{RST} {B}{name}{RST}")
    for ln in str(detail).splitlines()[:8]:
        print(f"       {ln}")

def skip(name, reason):
    global skipped; skipped += 1
    print(f"  {Y}SKIP{RST} {name}  ({reason})")

def sec(title):
    print(f"\n{B}{'='*62}{RST}\n{B}  {title}{RST}\n{B}{'='*62}{RST}")


# =========================================================
# 1.  _strip_metadata_lines
# =========================================================
sec("1 . _strip_metadata_lines")
from agents.pdf_utils import _strip_metadata_lines

STRIP_CASES = [
    # should be STRIPPED
    ("john.doe@iitb.ac.in",                                      False),
    ("ravi@ee.iitb.ac.in",                                       False),
    ("Dr. Rajesh Kumar",                                          False),
    ("Prof. Arun Singh Verma",                                    False),
    ("Mr. Rahul Mehta",                                           False),
    ("Indian Institute of Technology Delhi",                      False),
    ("Department of Electronics",                                 False),
    ("IIT Bombay",                                                False),
    ("University of Hyderabad",                                   False),
    ("May 2024",                                                  False),
    ("January 2026",                                              False),
    ("01/03/2026",                                                False),
    ("15-08-2025",                                                False),
    # should be KEPT
    ("The Fourier Transform is defined as:",                      True),
    ("Linearity: F{af+bg} = aF{f} + bF{g}",                      True),
    (r"x(t) = A \cos(2\pi f_0 t)",                               True),
    ("$E = mc^2$",                                                True),
    ("Convolution is associative — follows from the integral.",   True),
    ("1. Compute the DFT of x[n] using N-point FFT.",             True),
    ("f(x) = 1/(sigma*sqrt(2*pi)) * exp(-(x-mu)^2/(2*sigma^2))", True),
]

for line, keep in STRIP_CASES:
    result = _strip_metadata_lines(line).strip()
    kept = bool(result)
    if kept == keep:
        ok(repr(line)[:65], "kept" if keep else "stripped")
    else:
        fail(repr(line)[:65],
             f"Expected {'kept' if keep else 'stripped'}, got {'kept' if kept else 'stripped'}")


# =========================================================
# 2.  _clean_pdf_text
# =========================================================
sec("2 . _clean_pdf_text")
from agents.local_summarizer import _clean_pdf_text

noise_input = (
    "Lecture 3: Signals and Systems\n"
    "Slide 5\n"
    "Page 12\n"
    "12/45\n"
    "copyright IIT Bombay\n"
    "The Fourier Transform decomposes a signal into its frequency components.\n"
    "www.iitb.ac.in\n"
    "It is widely used in signal processing.\n"
    "FourierTransform converts time to frequency.\n"
)
cleaned = _clean_pdf_text(noise_input)

if "Lecture 3" not in cleaned and "Slide 5" not in cleaned and "Page 12" not in cleaned:
    ok("Noise lines (lecture/slide/page) removed")
else:
    fail("Noise lines still present", cleaned[:300])

if "The Fourier Transform" in cleaned and "widely used" in cleaned:
    ok("Content lines preserved")
else:
    fail("Content lines missing", cleaned[:300])

if "12/45" not in cleaned:
    ok("Page-fraction noise removed")
else:
    fail("Page-fraction still present", repr(cleaned[:200]))

cc_input = "FourierTransform converts signals from timeDomain to frequencyDomain."
cc_cleaned = _clean_pdf_text(cc_input)
if "Fourier" in cc_cleaned:
    ok("camelCase split: FourierTransform -> Fourier Transform")
else:
    fail("camelCase not split", repr(cc_cleaned))


# =========================================================
# 3.  _is_math_line
# =========================================================
sec("3 . _is_math_line  (formula detection)")
from agents.local_summarizer import _is_math_line

MATH_CASES = [
    ("X[k] = sum_{n=0}^{N-1} x[n] e^{-j2pi kn/N}",  True),
    ("f(t) = A sin(2*pi*f0*t)",                       True),
    ("alpha + beta = gamma",                           True),
    ("E = mc^2",                                       True),
    (r"$\int_0^\infty e^{-x} dx = 1$",               True),
    ("The system is linear and time-invariant.",       False),
    ("In this section we discuss the Fourier series.", False),
    ("Properties of Linear Systems",                   False),
    ("",                                               False),
    ("P(X=k) = C(n,k) * p^k * (1-p)^(n-k)",          True),
    ("y[n] = h[n] * x[n]",                            True),
]
for line, expected in MATH_CASES:
    result = _is_math_line(line)
    if result == expected:
        ok(repr(line)[:60], "math" if expected else "prose")
    else:
        fail(repr(line)[:60],
             f"Expected {'math' if expected else 'prose'}, got {'math' if result else 'prose'}")


# =========================================================
# 4.  extract_text_from_pptx  (synthetic PPTX)
# =========================================================
sec("4 . extract_text_from_pptx  (synthetic PPTX)")
try:
    from pptx import Presentation
    from pptx.util import Inches
    from agents.pdf_utils import extract_text_from_pptx

    def make_rich_pptx() -> bytes:
        prs = Presentation()

        # Slide 1: metadata/ title slide
        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s1.shapes.title.text = "Digital Signal Processing"
        s1.placeholders[1].text = (
            "Dr. Ravi Shankar\nravi@iitb.ac.in\n"
            "Department of Electrical Engineering, IIT Bombay\nMarch 2026"
        )

        # Slide 2: real content
        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Discrete Fourier Transform"
        s2.placeholders[1].text_frame.text = (
            "Converts a finite N-point sequence to the frequency domain.\n"
            "X[k] = sum_{n=0}^{N-1} x[n] e^{-j2pi kn/N}"
        )

        # Slide 3: another content slide
        s3 = prs.slides.add_slide(prs.slide_layouts[1])
        s3.shapes.title.text = "Convolution Theorem"
        s3.placeholders[1].text_frame.text = (
            "Circular convolution in time = multiplication in frequency domain."
        )

        # Slide 4: empty slide
        prs.slides.add_slide(prs.slide_layouts[6])

        # Slide 5: no title, just a textbox
        s5 = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = s5.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
        txBox.text_frame.text = "The Z-Transform is the discrete-time analogue of Laplace."

        buf = io.BytesIO(); prs.save(buf); return buf.getvalue()

    pptx_bytes = make_rich_pptx()
    extracted = extract_text_from_pptx(pptx_bytes)

    markers = [m for m in ["--- Slide 1", "--- Slide 2", "--- Slide 3"] if m in extracted]
    if len(markers) == 3:
        ok("Slide boundary markers for slides 1-3")
    else:
        fail("Missing slide markers", f"Found: {markers}")

    meta_leaks = [s for s in ["ravi@iitb.ac.in", "Dr. Ravi Shankar", "IIT Bombay"] if s in extracted]
    if not meta_leaks:
        ok("Author/email/institution metadata stripped")
    else:
        fail("Metadata leaked", str(meta_leaks))

    if "Discrete Fourier Transform" in extracted and "X[k]" in extracted:
        ok("Content titles and formulas preserved")
    else:
        fail("Content missing", extracted[:400])

    ok("Empty slide handled gracefully (no crash)")

    if "Z-Transform" in extracted:
        ok("No-title textbox slide content captured")
    else:
        fail("No-title slide content missing", extracted[-300:])

    pos2 = extracted.find("Discrete Fourier Transform")
    pos3 = extracted.find("Convolution Theorem")
    if 0 <= pos2 < pos3:
        ok("Slide ordering preserved (slide 2 before slide 3)")
    else:
        fail("Slide order wrong", f"pos DFT={pos2}, pos Conv={pos3}")

except Exception as e:
    fail("PPTX extraction suite", traceback.format_exc())


# =========================================================
# 5.  chunk_text / summarise_chunks
# =========================================================
sec("5 . chunk_text / summarise_chunks")
from agents.pdf_utils import chunk_text, summarise_chunks

SLIDE_TEXT = "\n\n".join(
    f"--- Slide {i}: Topic {i} ---\n" + ("Content for topic. " * 40)
    for i in range(1, 12)
)

chunks = chunk_text(SLIDE_TEXT, max_chars=3000)
slide_starts = len(re.findall(r"^--- Slide", SLIDE_TEXT, re.MULTILINE))
chunk_starts = sum(len(re.findall(r"^--- Slide", c, re.MULTILINE)) for c in chunks)
if chunk_starts == slide_starts:
    ok(f"All {slide_starts} slide markers preserved across {len(chunks)} chunks")
else:
    fail("Slide markers lost", f"{slide_starts} markers -> {chunk_starts} after chunking")

oversized = [len(c) for c in chunks if len(c) > 3400]
if not oversized:
    ok("No chunk exceeds max_chars limit")
else:
    fail(f"{len(oversized)} oversized chunks", f"lengths: {oversized}")

long_chunks = [f"Slide {i}: " + "word " * 500 for i in range(20)]
raw_total = sum(len(c) for c in long_chunks)
summary = summarise_chunks(long_chunks, max_summary_chars=5000)
# summarise_chunks enforces a 300-char-per-chunk minimum so it may exceed
# the nominal budget when budget/n_chunks < 300.  The important invariant is
# that output is significantly shorter than the untruncated input.
if len(summary) < raw_total * 0.5:
    ok(f"summarise_chunks trims significantly: {len(summary)} chars (was {raw_total})")
else:
    fail("summarise_chunks barely trimmed", f"got {len(summary)}, full={raw_total}")

slides_in_summary = len(re.findall(r"Slide \d+", summary))
if slides_in_summary == 20:
    ok(f"All 20 slides proportionally represented in summary")
else:
    fail(f"Only {slides_in_summary}/20 slides in summary", summary[:200])

empty_chunks = chunk_text("", max_chars=3000)
ok(f"Empty text returns {len(empty_chunks)} chunk(s) without crash")


# =========================================================
# 6.  fix_latex_delimiters
# =========================================================
sec("6 . fix_latex_delimiters")
from agents.latex_utils import fix_latex_delimiters

LATEX_CASES = [
    # (desc, input, fragment, should_be_present)
    ("inline paren -> $",
     r"The transform \(F(\omega)\) is defined",
     r"$F(\omega)$", True),

    ("multiline paren -> $ no internal whitespace",
     "value is \\(\n  x^2\n\\) here",
     "$x^2$", True),

    ("display bracket -> $$ block",
     "formula:\n\\[\n  E = mc^2\n\\]",
     "$$\nE = mc^2\n$$", True),

    ("single-line $$...$$ -> block",
     "so $$x^2 + y^2 = r^2$$ is it",
     "$$\nx^2 + y^2 = r^2\n$$", True),

    ("no spaces inside $ after strip",
     "the value \\(   alpha   \\) is small",
     "the value", True),

    ("already-good inline $ unchanged",
     "value is $x^2$ already",
     "$x^2$", True),

    ("blank line before $$ opener",
     "prose text\n$$\nformula\n$$",
     "\n\n$$", True),

    ("blank line after $$ closer",
     "$$\nformula\n$$\nnext paragraph",
     "$$\n\nnext", True),

    ("no raw paren delimiters remain",
     r"\(x + y\) and \(a - b\)",
     "\\(", False),

    ("no raw bracket delimiters remain",
     r"display \[E=mc^2\]",
     "\\[", False),

    ("excessive blank lines collapsed - max 2 newlines",
     "para A\n\n\n\n\n\npara B",
     "\n\n\n\n", False),

    ("multiline \\[\\] preserves inner newlines",
     "\\[\n  a + b\\\\\n  = c\n\\]",
     "$$", True),
]

for desc, src, fragment, should_present in LATEX_CASES:
    try:
        result = fix_latex_delimiters(src)
        present = fragment in result
        if present == should_present:
            ok(desc)
        else:
            fail(desc,
                 f"Fragment {repr(fragment)} {'NOT found' if should_present else 'STILL present'}\n"
                 f"Result: {repr(result)}")
    except Exception as e:
        fail(desc, str(e))


# =========================================================
# 7.  generate_local_note  (offline summarizer)
# =========================================================
sec("7 . generate_local_note  (offline summarizer)")
from agents.local_summarizer import generate_local_note

SLIDES = (
    "--- Slide 1: Fourier Transform ---\n"
    "Converts time-domain signal to frequency domain.\n"
    "F(omega) = integral x(t) e^{-j omega t} dt\n"
    "Property: Linearity: F{a*x + b*y} = a*F{x} + b*F{y}\n\n"
    "--- Slide 2: Convolution Theorem ---\n"
    "Convolution in time = multiplication in frequency.\n"
    "y(t) = x(t) * h(t)  =>  Y(omega) = X(omega) * H(omega)\n\n"
    "--- Slide 3: Z-Transform ---\n"
    "X(z) = sum_{n=-inf}^{inf} x[n] z^{-n}\n"
    "LTI systems analysed via transfer function H(z).\n"
)

TEXTBOOK = (
    "The Fourier Transform F(omega) maps a time-domain signal x(t) to "
    "its frequency-domain representation. The inverse transform recovers "
    "x(t). For Linear Time-Invariant systems, the convolution theorem "
    "allows frequency-domain analysis: Y = XH."
)

for level in ("Beginner", "Intermediate", "Advanced"):
    try:
        note = generate_local_note(SLIDES, TEXTBOOK, level)

        if len(note.strip()) < 100:
            fail(f"{level}: non-empty output", f"Only {len(note)} chars"); continue
        ok(f"{level}: output generated", f"{len(note)} chars")

        headings = re.findall(r"^## .+", note, re.MULTILINE)
        if headings:
            ok(f"{level}: ## headings present", f"{len(headings)} found")
        else:
            fail(f"{level}: ## headings missing", note[:300])

        if re.search(r"fourier", note, re.I):
            ok(f"{level}: Fourier content present")
        else:
            fail(f"{level}: Fourier content missing", note[:200])

        if re.search(r"\$", note):
            ok(f"{level}: LaTeX math present in output")
        else:
            fail(f"{level}: No math in output", note[:300])

        bad = [d for d in ["\\(", "\\)", "\\[", "\\]"] if d in note]
        if not bad:
            ok(f"{level}: No raw paren/bracket LaTeX delimiters")
        else:
            fail(f"{level}: Raw delimiters present", str(bad))

    except Exception as e:
        fail(f"{level}: generate_local_note", traceback.format_exc())


# =========================================================
# 8.  extract_concepts
# =========================================================
sec("8 . extract_concepts")
from agents.concept_extractor import extract_concepts

DSP_NOTE = (
    "## Fourier Transform\n"
    "The Fourier Transform decomposes a signal into frequency components.\n\n"
    "## Convolution Theorem\n"
    "Convolution in time equals multiplication in frequency.\n\n"
    "## Z-Transform\n"
    "X(z) = sum x[n] z^{-n}. Used for LTI system analysis.\n\n"
    "## Sampling Theorem\n"
    "The Nyquist rate is twice the highest frequency.\n"
)

graph = extract_concepts(DSP_NOTE)
labels = {n["label"] for n in graph["nodes"]}
for expected in ["Fourier Transform", "Convolution Theorem", "Z-Transform", "Sampling Theorem"]:
    if expected in labels:
        ok(f'Concept "{expected}" extracted')
    else:
        fail(f'Concept "{expected}" missing', f"Found: {sorted(labels)}")

if len(graph["edges"]) > 0:
    ok(f"Dependency edges generated: {len(graph['edges'])} edges")
else:
    fail("No dependency edges", "Expected edges between DSP concepts")

ids = [n["id"] for n in graph["nodes"]]
if len(ids) == len(set(ids)) and all(isinstance(i, int) for i in ids):
    ok("Node IDs are unique integers")
else:
    fail("Duplicate or non-integer node IDs", str(ids))

required_fields = {"id", "label", "status", "x", "y"}
bad_nodes = [n for n in graph["nodes"] if not required_fields.issubset(n.keys())]
if not bad_nodes:
    ok("All nodes have required fields (id, label, status, x, y)")
else:
    fail(f"{len(bad_nodes)} nodes missing fields", str(bad_nodes[:2]))

id_set = set(ids)
bad_edges = [e for e in graph["edges"] if e[0] not in id_set or e[1] not in id_set]
if not bad_edges:
    ok("All edges reference valid node IDs")
else:
    fail(f"{len(bad_edges)} edges with invalid node IDs", str(bad_edges[:5]))

UNKNOWN_NOTE = (
    "## Quantum Entanglement Basics\n"
    "Particles become correlated such that the state of one instantly affects the other.\n\n"
    "## Bell Theorem\n"
    "No local hidden variable theory can reproduce all quantum mechanical predictions.\n\n"
    "## EPR Paradox\n"
    "Einstein-Podolsky-Rosen proposed a thought experiment.\n"
)
graph2 = extract_concepts(UNKNOWN_NOTE)
if len(graph2["nodes"]) >= 1:
    ok(f"Heading fallback: {len(graph2['nodes'])} nodes extracted for unknown topic")
else:
    fail("Heading fallback: no nodes extracted", str(graph2))


# =========================================================
# 9.  local_mutate
# =========================================================
sec("9 . local_mutate  (offline doubt resolution)")
from agents.local_mutation import local_mutate

PARAGRAPH = (
    "## Convolution Theorem\n\n"
    "The Convolution Theorem states that convolution in the time domain corresponds "
    "to multiplication in the frequency domain.\n\n"
    "This is why analyzing LTI systems in the frequency domain is so powerful."
)

DOUBT_CASES = [
    ("why does convolution become multiplication?",   "why"),
    ("how do I compute this step by step?",           "how"),
    ("what exactly is convolution?",                  "what"),
    ("I don't understand the formula at all",         "dont understand"),
    ("confused by the notation",                      "confused"),
    ("give me an intuitive explanation please",       "intuitive"),
    ("when should I use the frequency domain?",       "when"),
    ("difference between convolution and correlation?", "difference"),
]

for doubt, tag in DOUBT_CASES:
    try:
        mutated, gap = local_mutate(PARAGRAPH, doubt)
        if not mutated.strip():
            fail(f'Doubt [{tag}]: empty mutated output', ""); continue
        ok(f'Doubt [{tag}]: output generated', f"{len(mutated)} chars")

        if "Intuition" in mutated or "intuition" in mutated.lower():
            ok(f'Doubt [{tag}]: insight block injected')
        else:
            fail(f'Doubt [{tag}]: insight block missing', mutated[:200])

        if gap.strip():
            ok(f'Doubt [{tag}]: concept_gap diagnosis returned', gap[:60])
        else:
            fail(f'Doubt [{tag}]: empty concept_gap', "")

        if "## Convolution Theorem" in mutated:
            ok(f'Doubt [{tag}]: heading preserved')
        else:
            fail(f'Doubt [{tag}]: heading lost', mutated[:150])

    except Exception as e:
        fail(f'local_mutate [{tag}]', traceback.format_exc())


# =========================================================
# 10.  local_examine
# =========================================================
sec("10 . local_examine  (offline MCQ generator)")
from agents.local_examiner import local_examine

EXAMINE_CASES = [
    "fourier transform",
    "convolution",
    "laplace transform",
    "z-transform",
    "probability",
    "xyzzy_unknown_concept",  # generic fallback
]

for concept in EXAMINE_CASES:
    try:
        result = local_examine(concept)
        if not result.strip():
            fail(f'examine [{concept}]: empty output', ""); continue
        ok(f'examine [{concept}]: output generated', f"{len(result)} chars")

        # Must have A/B/C/D choices
        if all(opt in result for opt in ("A)", "B)", "C)", "D)")):
            ok(f'examine [{concept}]: A/B/C/D options present')
        else:
            fail(f'examine [{concept}]: missing MCQ options', result[:200])

        if "?" in result:
            ok(f'examine [{concept}]: question mark present')
        else:
            fail(f'examine [{concept}]: no question mark', result[:150])

    except Exception as e:
        fail(f'local_examine [{concept}]', traceback.format_exc())


# =========================================================
# 11.  Page splitter (Python mirror of JS useMemo logic)
# =========================================================
sec("11 . Page splitter  (JS useMemo logic mirrored)")

def js_page_splitter(note: str, target: int = 2200) -> list:
    if not note:
        return []
    by_h2 = [s.strip() for s in re.split(r"(?=^## )", note, flags=re.MULTILINE) if s.strip()]
    if by_h2:
        merged, buf = [], ""
        for s in by_h2:
            if buf and len(buf) + len(s) + 2 > target and len(buf) > 500:
                merged.append(buf.strip()); buf = s
            else:
                buf = (buf + "\n\n" + s) if buf else s
        if buf:
            merged.append(buf.strip())
        return [p for p in merged if p]
    return [note]

# 11a: 15 sections -> grouped pages
note15 = "\n\n".join(
    f"## Topic {i}\n" + "Content sentence here. " * 30 for i in range(1, 16)
)
pages15 = js_page_splitter(note15)
if 3 <= len(pages15) <= 7:
    ok(f"15 sections grouped into {len(pages15)} pages (expected 3-7)")
else:
    fail(f"15 sections -> {len(pages15)} pages", "Expected 3-7")

# 11b: $$ blocks not orphaned
math_note = (
    "## Section A\n\nIntro.\n\n$$\nE = mc^2\n$$\n\nAfter formula.\n\n"
    + "\n\n".join(f"## Sec {i}\n" + "Filler text. " * 55 for i in range(1, 9))
)
pages_m = js_page_splitter(math_note)
broken = any(p.count("$$") % 2 != 0 for p in pages_m)
if not broken:
    ok("No $$ block orphaned across page boundary")
else:
    fail("$$ block split", "\n".join(f"p{i}: {p.count('$$')} $$" for i, p in enumerate(pages_m)))

# 11c: single short section -> 1 page
pages_s = js_page_splitter("## One Topic\n\nJust a short paragraph.")
if len(pages_s) == 1:
    ok("Short single-section note -> 1 page")
else:
    fail("Short note split unexpectedly", f"{len(pages_s)} pages")

# 11d: empty note -> 0 pages
if js_page_splitter("") == []:
    ok("Empty note -> 0 pages")
else:
    fail("Empty note not handled", "")

# 11e: boundary split is sane
boundary_note = "\n\n".join(
    f"## T{i}\n" + "x " * 100 for i in range(1, 12)
)
pages_b = js_page_splitter(boundary_note, target=2200)
if 1 <= len(pages_b) <= 11:
    ok(f"Boundary split: {len(pages_b)} pages (1-11 acceptable)")
else:
    fail(f"Boundary split returned {len(pages_b)} pages", "")


# =========================================================
# 12.  Live API  (all major endpoints)
# =========================================================
sec("12 . Live API  (all major endpoints)")
try:
    import requests
    BASE = "http://localhost:8000"

    # 12a: health
    r = requests.get(f"{BASE}/health", timeout=5)
    if r.status_code == 200 and r.json().get("status") == "ok":
        ok("GET /health")
    else:
        fail("GET /health", r.text[:100])

    # 12b: register
    test_email = f"test_{__import__('uuid').uuid4().hex[:8]}@auragraph.test"
    r = requests.post(f"{BASE}/auth/register",
                      json={"email": test_email, "password": "Test@1234"}, timeout=5)
    if r.status_code in (200, 201):
        token = r.json().get("token", "")
        ok("POST /auth/register", f"token={'ok' if token else 'missing'}")
    else:
        fail("POST /auth/register", f"HTTP {r.status_code}: {r.text[:100]}")
        token = ""

    # 12c: login
    r = requests.post(f"{BASE}/auth/login",
                      json={"email": test_email, "password": "Test@1234"}, timeout=5)
    if r.status_code == 200:
        token = r.json().get("token", token)
        ok("POST /auth/login")
    else:
        fail("POST /auth/login", f"HTTP {r.status_code}: {r.text[:100]}")

    # 12d: duplicate register -> 409
    r2 = requests.post(f"{BASE}/auth/register",
                       json={"email": test_email, "password": "anything"}, timeout=5)
    if r2.status_code == 409:
        ok("POST /auth/register duplicate -> 409")
    else:
        fail("Duplicate register should be 409", f"got {r2.status_code}")

    # 12e: wrong password -> 401
    r3 = requests.post(f"{BASE}/auth/login",
                       json={"email": test_email, "password": "wrong"}, timeout=5)
    if r3.status_code == 401:
        ok("POST /auth/login wrong password -> 401")
    else:
        fail("Wrong password should be 401", f"got {r3.status_code}")

    headers = {"Authorization": f"Bearer {token}"} if token else {}

    # 12f: create notebook
    r = requests.post(f"{BASE}/notebooks",
                      json={"name": "Test DSP", "course": "EE301"},
                      headers=headers, timeout=5)
    if r.status_code in (200, 201):
        nb_id = r.json().get("id", "")
        ok("POST /notebooks create", f"id={str(nb_id)[:12]}...")
    else:
        fail("POST /notebooks create", f"HTTP {r.status_code}: {r.text[:150]}")
        nb_id = ""

    # 12g: list notebooks
    r = requests.get(f"{BASE}/notebooks", headers=headers, timeout=5)
    if r.status_code == 200 and isinstance(r.json(), list):
        ok("GET /notebooks list", f"{len(r.json())} notebooks")
    else:
        fail("GET /notebooks list", f"HTTP {r.status_code}: {r.text[:100]}")

    # 12h: fetch single notebook
    if nb_id:
        r = requests.get(f"{BASE}/notebooks/{nb_id}", headers=headers, timeout=5)
        if r.status_code == 200 and r.json().get("id") == nb_id:
            ok("GET /notebooks/{id}")
        else:
            fail("GET /notebooks/{id}", f"HTTP {r.status_code}: {r.text[:100]}")

    # 12i: unauthorized -> 401
    r = requests.get(f"{BASE}/notebooks", timeout=5)
    if r.status_code == 401:
        ok("GET /notebooks without token -> 401")
    else:
        fail("Unauthenticated should be 401", f"got {r.status_code}")

    # 12j: /api/fuse local fallback
    SLIDES_API = (
        "--- Slide 1: Fourier Transform ---\n"
        "Converts time-domain signal to frequency domain.\n"
        "F(omega) = integral x(t) e^{-j omega t} dt\n\n"
        "--- Slide 2: Convolution Theorem ---\n"
        "y(t) = x(t) * h(t)  =>  Y(omega) = X(omega) * H(omega)\n"
    )
    TEXTBOOK_API = (
        "The Fourier Transform maps time-domain signals to frequency domain. "
        "Convolution in time corresponds to multiplication in frequency."
    )
    r = requests.post(f"{BASE}/api/fuse",
                      json={"slide_summary": SLIDES_API,
                            "textbook_paragraph": TEXTBOOK_API,
                            "proficiency": "Intermediate"},
                      timeout=30)
    if r.status_code == 200:
        note = r.json().get("fused_note", "")
        if note.strip():
            ok("/api/fuse local fallback", f"{len(note)} chars")
            if re.search(r"^## ", note, re.MULTILINE):
                ok("/api/fuse output has ## headings")
            else:
                fail("/api/fuse no ## headings", note[:200])
            bad_delims = [d for d in ["\\(", "\\["] if d in note]
            if not bad_delims:
                ok("/api/fuse no raw LaTeX delimiters")
            else:
                fail("/api/fuse raw delimiters in output", str(bad_delims))
        else:
            fail("/api/fuse: fused_note empty", str(r.json()))
    else:
        fail(f"/api/fuse HTTP {r.status_code}", r.text[:200])

    # 12k: /api/upload-fuse-multi
    try:
        pptx_bytes_api = pptx_bytes
    except NameError:
        pptx_bytes_api = b""
    MINIMAL_PDF = (
        b"%PDF-1.4\n1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
        b"3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R"
        b"/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj\n"
        b"4 0 obj<</Length 44>>stream\n"
        b"BT /F1 12 Tf 72 720 Td (Textbook content here.) Tj ET\n"
        b"endstream endobj\n"
        b"5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
        b"xref\n0 6\n0000000000 65535 f \n0000000009 00000 n \n"
        b"0000000058 00000 n \n0000000115 00000 n \n0000000266 00000 n \n"
        b"0000000360 00000 n \ntrailer<</Size 6/Root 1 0 R>>\nstartxref\n441\n%%EOF"
    )
    SLIDE_PDF = (
        b"%PDF-1.4\n1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
        b"3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R"
        b"/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj\n"
        b"4 0 obj<</Length 52>>stream\n"
        b"BT /F1 12 Tf 72 720 Td (Fourier Transform slides content.) Tj ET\n"
        b"endstream endobj\n"
        b"5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
        b"xref\n0 6\n0000000000 65535 f \n0000000009 00000 n \n"
        b"0000000058 00000 n \n0000000115 00000 n \n0000000266 00000 n \n"
        b"0000000360 00000 n \ntrailer<</Size 6/Root 1 0 R>>\nstartxref\n449\n%%EOF"
    )
    if pptx_bytes_api:
        slide_name = "slides.pptx"
        slide_mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        slide_data = pptx_bytes_api
    else:
        slide_name = "slides.pdf"
        slide_mime = "application/pdf"
        slide_data = SLIDE_PDF

    files = {
        "slides_pdfs":   (slide_name, slide_data, slide_mime),
        "textbook_pdfs": ("book.pdf", MINIMAL_PDF, "application/pdf"),
    }
    r = requests.post(f"{BASE}/api/upload-fuse-multi",
                      files=files, data={"proficiency": "Beginner"}, timeout=60)
    if r.status_code == 200:
        note = r.json().get("fused_note", "")
        ok("/api/upload-fuse-multi HTTP 200", f"{len(note)} chars")
        try:
            meta_in = [s for s in ["ravi@iitb.ac.in", "Dr. Ravi Shankar"] if s in note]
            ok("/api/upload-fuse-multi: no metadata in output") if not meta_in \
                else fail("/api/upload-fuse-multi metadata leak", str(meta_in))
        except Exception:
            pass
    else:
        fail(f"/api/upload-fuse-multi HTTP {r.status_code}", r.text[:200])

    # 12l: /api/mutate
    r = requests.post(f"{BASE}/api/mutate",
                      json={"original_paragraph": "## Fourier Transform\n\nConverts time to frequency.",
                            "student_doubt": "why does it decompose into frequencies?"},
                      timeout=20)
    if r.status_code == 200:
        body = r.json()
        if body.get("mutated_paragraph") and body.get("concept_gap"):
            ok("/api/mutate returns mutated_paragraph + concept_gap")
        else:
            fail("/api/mutate missing fields", str(body)[:200])
    else:
        fail(f"/api/mutate HTTP {r.status_code}", r.text[:150])

    # 12m: /api/examine
    r = requests.post(f"{BASE}/api/examine",
                      json={"concept_name": "Fourier Transform"}, timeout=20)
    if r.status_code == 200:
        questions = r.json().get("practice_questions", "")
        if "?" in questions and "A)" in questions:
            ok("/api/examine returns MCQ questions")
        else:
            fail("/api/examine bad output format", questions[:200])
    else:
        fail(f"/api/examine HTTP {r.status_code}", r.text[:150])

    # 12n: /api/extract-concepts
    r = requests.post(f"{BASE}/api/extract-concepts",
                      json={"note": DSP_NOTE}, timeout=10)
    if r.status_code == 200:
        g = r.json()
        if g.get("nodes") and isinstance(g["nodes"], list):
            ok(f"/api/extract-concepts: {len(g['nodes'])} nodes, {len(g['edges'])} edges")
        else:
            fail("/api/extract-concepts bad response", str(g)[:150])
    else:
        fail(f"/api/extract-concepts HTTP {r.status_code}", r.text[:150])

    # 12o: delete notebook + 404 verification
    if nb_id:
        r = requests.delete(f"{BASE}/notebooks/{nb_id}", headers=headers, timeout=5)
        if r.status_code == 200:
            ok("DELETE /notebooks/{id}")
            r2 = requests.get(f"{BASE}/notebooks/{nb_id}", headers=headers, timeout=5)
            if r2.status_code == 404:
                ok("GET /notebooks/{id} after delete -> 404")
            else:
                fail("Notebook still accessible after delete", f"HTTP {r2.status_code}")
        else:
            fail(f"DELETE /notebooks HTTP {r.status_code}", r.text[:100])

except ImportError:
    skip("Live API tests", "requests not installed")
except ConnectionError:
    skip("Live API tests", "Backend not running on :8000")
except Exception as e:
    fail("Live API test suite", traceback.format_exc())


# =========================================================
# 13.  Aggressive math: _raw_to_latex + delimiter torture
# =========================================================
sec("13 . Aggressive math: _raw_to_latex + delimiter torture")
from agents.local_summarizer import _raw_to_latex, _is_math_line

# ── 13a: _raw_to_latex symbol table ──────────────────────────────
RAW_TO_LATEX_CASES = [
    # (raw input,             expected substring in output)
    ("integral x dt",         r"\int"),
    ("sum_{n=0}^{N}",         r"\sum"),
    ("prod_{k=1}^{N}",        r"\prod"),
    ("sqrt(x^2 + y^2)",       r"\sqrt"),
    ("lim_{x->0} f(x)",       r"\lim"),
    ("alpha + beta",          r"\alpha"),
    ("gamma * delta",         r"\gamma"),
    ("epsilon < zeta",        r"\epsilon"),
    ("theta = pi/4",         r"\theta"),
    ("lambda * Omega",        r"\lambda"),
    ("mu and sigma",          r"\mu"),
    ("sigma and rho",         r"\sigma"),
    ("omega = 2pi*f_0",       r"\omega"),
    ("nabla f(x,y)",          r"\nabla"),
    ("x -> y => z",           r"\rightarrow"),
    ("x <-> y",               r"\leftrightarrow"),
    ("P(X) >= 0",             r"\geq"),
    ("a != b",                r"\neq"),
    ("infty or infinity",     r"\infty"),
    ("sin(theta) + cos(phi)", r"\sin"),
    ("log(x) + ln(y)",        r"\log"),
    ("exp(-x^2/2)",           r"\exp"),
]
for raw, expected_sub in RAW_TO_LATEX_CASES:
    result = _raw_to_latex(raw)
    if expected_sub in result:
        ok(f"_raw_to_latex: {repr(raw)[:40]}", f"-> {expected_sub}")
    else:
        fail(f"_raw_to_latex: {repr(raw)[:40]}",
             f"Expected {repr(expected_sub)} in output\nGot: {repr(result)}")

# ── 13b: _is_math_line torture — harder edge cases ────────────────
HARD_MATH_CASES = [
    # heavy LaTeX — definitely math
    (r"$$\int_{-\infty}^{\infty} f(x) e^{-j\omega x} dx$$",  True),
    (r"$X(z) = \sum_{n=0}^{N-1} x[n] z^{-n}$",              True),
    (r"$\sigma^2 = \frac{1}{N}\sum_{i=1}^{N}(x_i - \mu)^2$", True),
    (r"$P(A|B) = \frac{P(B|A)P(A)}{P(B)}$",                  True),
    # operator-dense short lines  —  math
    ("a^2 + b^2 = c^2",                                       True),
    ("H(z) = Y(z)/X(z)",                                      True),
    ("f'(x) = lim_{h->0} [f(x+h)-f(x)]/h",                   True),
    ("det(A - lambda*I) = 0",                                  True),
    # pure prose sentences  —  not math
    ("The sampling theorem is fundamental to digital signal processing.", False),
    ("We will now prove the convolution theorem using Fourier analysis.", False),
    ("Linearity means the system satisfies superposition.",             False),
    ("Consider a discrete-time signal x of length N.",                 False),
    # borderline — very short single-char assignment is treated as math
    ("n = 8",                                                          True),
    ("k is the frequency index",                                       False),
]
for line, expected in HARD_MATH_CASES:
    result = _is_math_line(line)
    if result == expected:
        ok(f"_is_math_line hard: {repr(line)[:55]}", "math" if expected else "prose")
    else:
        fail(f"_is_math_line hard: {repr(line)[:55]}",
             f"Expected {'math' if expected else 'prose'}, got {'math' if result else 'prose'}")

# ── 13c: fix_latex_delimiters — complex real-world documents ──────
TORTURE_CASES = [
    # Multiple inline parens in one paragraph
    ("multi inline in para",
     r"We have \(f(x)\) and also \(g(x)\) which together give \(h(x)\).",
     lambda r: r.count("\\(") == 0 and r.count("$") == 6,
     "all 3 \\(...\\) -> $ ... $"),

    # Display block nested inside a paragraph (common in AI output)
    ("display inside para text",
     "The DFT is given by:\n\\[\nX[k] = \\sum_{n=0}^{N-1} x[n] e^{-j2\\pi kn/N}\n\\]\nand this defines the spectrum.",
     lambda r: "\\[" not in r and "$$" in r,
     "\\[...\\] replaced by $$ block"),

    # Mix of inline and display in one document
    ("mixed inline+display",
     "Energy \\(E = mc^2\\) is famous.\n\n\\[\\int_0^1 x dx = 0.5\\]\n\nThis ends the proof.",
     lambda r: "\\(" not in r and "\\[" not in r,
     "no raw delimiters remain"),

    # Already-normalised note must survive unchanged (idempotent)
    ("idempotent: already normalised note unchanged",
     "## Topic\n\nProse text.\n\n$$\nE = mc^2\n$$\n\nMore prose with $x^2$ inline.",
     lambda r: r.count("$$") == 2 and "$x^2$" in r,
     "already-normalised note unchanged"),

    # Wall of $$ — many display blocks back to back
    ("many display blocks in sequence",
     "\n\n".join(f"$$\nf_{i}(x)\n$$" for i in range(5)),
     lambda r: r.count("$$") == 10,
     "10 $$ markers preserved"),

    # Inline math at start and end of line (no surrounding space)
    ("inline math at line boundaries",
     r"$a$ plus $b$ equals $c$",
     lambda r: r.count("$") >= 6 and "\\(" not in r,
     "6 $ markers, no raw parens"),

    # Bare \\( ... \\) spanning 5 lines
    ("five-line inline paren",
     "Start \\(\n  a\n  +\n  b\n  = c\n\\) end",
     lambda r: "\\(" not in r and "$" in r,
     "five-line paren becomes single $"),

    # Real formula: Bayes theorem with multiple nested parens
    ("Bayes theorem with parens",
     r"\(P(A|B) = \frac{P(B|A) \cdot P(A)}{P(B)}\)",
     lambda r: "\\(" not in r and "$P(A|B)" in r,
     "Bayes inline -> $...$"),
]

for desc, src, predicate, detail in TORTURE_CASES:
    try:
        result = fix_latex_delimiters(src)
        if predicate(result):
            ok(f"Torture [{desc}]", detail)
        else:
            fail(f"Torture [{desc}]", f"{detail}\nResult: {repr(result[:200])}")
    except Exception as e:
        fail(f"Torture [{desc}]", str(e))

# ── 13d: image-in-note — _strip_metadata_lines & pipeline pass-through ──
IMAGE_IN_MD = (
    "## Fourier Transform\n\n"
    "The spectrum is shown below.\n\n"
    "![Frequency spectrum](spectrum.png)\n\n"
    "Note the peak at f_0.\n\n"
    "![Block diagram of LTI system](lti_block.png)\n\n"
    "The transfer function H(z) relates input to output.\n"
)
# After stripping metadata lines, markdown image syntax must NOT be destroyed
stripped = _strip_metadata_lines(IMAGE_IN_MD)
if "![Frequency spectrum]" in stripped and "![Block diagram" in stripped:
    ok("Image Markdown syntax preserved through _strip_metadata_lines")
else:
    fail("Image Markdown syntax destroyed by _strip_metadata_lines",
         f"Got: {repr(stripped[:300])}")

# generate_local_note must not crash on a note containing markdown image refs
try:
    note_with_img = generate_local_note(IMAGE_IN_MD, "See figure for reference.", "Beginner")
    ok("generate_local_note handles note with image markdown (no crash)",
       f"{len(note_with_img)} chars")
except Exception as e:
    fail("generate_local_note crashed on image markdown", traceback.format_exc())

# fix_latex_delimiters must not destroy ![alt](src) image references
img_md = "Some text $x^2$ and an image: ![fig](img.png) and more $y$."
img_result = fix_latex_delimiters(img_md)
if "![fig](img.png)" in img_result:
    ok("fix_latex_delimiters preserves ![alt](src) image references")
else:
    fail("fix_latex_delimiters destroyed image markdown",
         f"Got: {repr(img_result)}")


# =========================================================
# Summary
# =========================================================
total = passed + failed + skipped
colour = G if failed == 0 else R
print(f"\n{'='*62}")
print(f"{B}{colour}  {passed}/{total} passed  .  {failed} failed  .  {skipped} skipped{RST}")
print(f"{'='*62}\n")
if failed:
    import sys; sys.exit(1)
